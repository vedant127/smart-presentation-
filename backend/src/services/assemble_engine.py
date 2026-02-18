import sys
import os
import json
import re
import uuid
import zipfile
import shutil
import tempfile
from datetime import datetime

try:
    from lxml import etree
    from pptx import Presentation
    from pptx.util import Inches, Emu
except ImportError:
    print(json.dumps({"success": False, "error": "python-pptx or lxml not installed. Run: pip install python-pptx lxml"}))
    sys.exit(1)


# ─────────────────────────────────────────────────────────────
# CONSTANTS
# ─────────────────────────────────────────────────────────────
SLIDE_REL_TYPE = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide'
SLIDE_LAYOUT_REL_TYPE = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout'
NS_P = 'http://schemas.openxmlformats.org/presentationml/2006/main'
NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
NS_PR = 'http://schemas.openxmlformats.org/package/2006/relationships'
NS_CT = 'http://schemas.openxmlformats.org/package/2006/content-types'


# ─────────────────────────────────────────────────────────────
# HELPER: Parse .rels XML from a ZIP
# ─────────────────────────────────────────────────────────────
def parse_rels(zip_file, rels_path):
    try:
        data = zip_file.read(rels_path)
        root = etree.fromstring(data)
        rels = []
        for rel in root:
            rels.append({
                'id': rel.get('Id'),
                'type': rel.get('Type'),
                'target': rel.get('Target'),
                'mode': rel.get('TargetMode', 'Internal')
            })
        return rels
    except KeyError:
        return []


# ─────────────────────────────────────────────────────────────
# HELPER: Build .rels XML bytes
# ─────────────────────────────────────────────────────────────
def build_rels_xml(rels):
    lines = [
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>',
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
    ]
    for rel in rels:
        rid = rel['id']
        rtype = rel['type']
        rtarget = rel['target']
        mode_attr = ''
        if rel.get('mode') == 'External':
            mode_attr = ' TargetMode="External"'
        lines.append(
            '  <Relationship Id="{}" Type="{}" Target="{}"{} />'.format(rid, rtype, rtarget, mode_attr)
        )
    lines.append('</Relationships>')
    return '\n'.join(lines).encode('utf-8')


# ─────────────────────────────────────────────────────────────
# HELPER: Normalize a string for fuzzy filename matching
# ─────────────────────────────────────────────────────────────
def normalize_for_match(s):
    """Lowercase, strip, collapse spaces, remove special chars for comparison."""
    s = s.lower().strip()
    s = re.sub(r'[\s_+\-]+', '_', s)  # normalize separators to underscore
    s = re.sub(r'[^a-z0-9_]', '', s)  # remove anything else
    return s


# ─────────────────────────────────────────────────────────────
# HELPER: Build all filename variants from plot criteria
# Supports both naming conventions:
#   underscore:  dubai_residential_apartments_luxury.pptx
#   plus-space:  dubai + residential + apartments + luxury.pptx
# ─────────────────────────────────────────────────────────────
def build_varying_filenames(criteria_values):
    """Returns list of candidate filenames to try, in priority order."""
    parts = [v.strip().lower() for v in criteria_values if v and v.strip()]
    if not parts:
        return []
    return [
        '_'.join(parts) + '.pptx',          # dubai_residential_apartments_luxury.pptx
        ' + '.join(parts) + '.pptx',         # dubai + residential + apartments + luxury.pptx
    ]


def build_varying_filename(criteria_values):
    """Legacy compat: returns the underscore variant."""
    parts = [v.strip().lower() for v in criteria_values if v and v.strip()]
    if not parts:
        return None
    return '_'.join(parts) + '.pptx'


# ─────────────────────────────────────────────────────────────
# HELPER: Find a PPTX file in a folder
# Tries multiple naming conventions (underscore, plus-space)
# Falls back to fuzzy normalized match
# ─────────────────────────────────────────────────────────────
def find_pptx_in_folder(folder_path, filename=None):
    if not os.path.isdir(folder_path):
        return None

    all_files = [f for f in os.listdir(folder_path)
                 if f.lower().endswith('.pptx') and not f.startswith('~$')]

    if not all_files:
        return None

    # Non-varying: if no filename specified, return first file found
    if filename is None:
        return os.path.join(folder_path, all_files[0])

    # 1. Exact case-insensitive match
    filename_lower = filename.lower()
    for f in all_files:
        if f.lower() == filename_lower:
            return os.path.join(folder_path, f)

    # 2. Fuzzy normalized match (ignores separator style)
    target_norm = normalize_for_match(os.path.splitext(filename)[0])
    for f in all_files:
        f_norm = normalize_for_match(os.path.splitext(f)[0])
        if f_norm == target_norm:
            return os.path.join(folder_path, f)

    return None


def find_pptx_for_criteria(folder_path, criteria_values, partial_fallback=True):
    """
    Try all filename variants for the given criteria values.
    Also tries partial match (City + AssetType only) as fallback.
    Returns (path, matched_filename) or (None, None).
    """
    # Try all full variants
    for fname in build_varying_filenames(criteria_values):
        path = find_pptx_in_folder(folder_path, fname)
        if path:
            return path, fname

    # Try partial match (first 2 criteria: City + AssetType)
    if partial_fallback and len(criteria_values) > 2:
        partial_values = criteria_values[:2]
        for fname in build_varying_filenames(partial_values):
            path = find_pptx_in_folder(folder_path, fname)
            if path:
                return path, fname + ' [partial]'

    return None, None


# ─────────────────────────────────────────────────────────────
# HELPER: Deduplicate plot combinations
# ─────────────────────────────────────────────────────────────
def deduplicate_plots(plots, criteria_keys):
    seen = set()
    unique = []
    for plot in plots:
        criteria = plot.get('criteria', {})
        key_parts = []
        for k in criteria_keys:
            val = next((v for ck, v in criteria.items() if ck.lower() == k.lower()), '')
            key_parts.append(str(val).strip().lower())
        key = tuple(key_parts)
        if key not in seen:
            seen.add(key)
            unique.append(criteria)
    return unique


# ─────────────────────────────────────────────────────────────
# CORE: PptxAssembler
# Assembles a PPTX by copying slides at the ZIP level.
# This is the ONLY reliable way to copy slides with images/charts.
# ─────────────────────────────────────────────────────────────
class PptxAssembler:

    def __init__(self, template_path=None):
        self.tmp_path = tempfile.mktemp(suffix='.pptx')
        self.slide_count = 0

        if template_path and os.path.exists(template_path):
            shutil.copy2(template_path, self.tmp_path)
            self._clear_slides()
        else:
            prs = Presentation()
            prs.save(self.tmp_path)
            self._clear_slides()

    def _clear_slides(self):
        """Remove all slides from the PPTX, keeping master/layouts."""
        tmp_out = self.tmp_path + '.clearing'

        with zipfile.ZipFile(self.tmp_path, 'r') as zin:
            names = zin.namelist()
            slide_files = set(n for n in names if re.match(r'ppt/slides/slide\d+\.xml$', n))
            slide_rels = set(n for n in names if re.match(r'ppt/slides/_rels/slide\d+\.xml\.rels$', n))
            to_remove = slide_files | slide_rels

            # Read and modify presentation.xml
            prs_xml = zin.read('ppt/presentation.xml')
            prs_root = etree.fromstring(prs_xml)
            sldIdLst = prs_root.find('.//{%s}sldIdLst' % NS_P)
            if sldIdLst is not None:
                for child in list(sldIdLst):
                    sldIdLst.remove(child)
            mod_prs_xml = etree.tostring(prs_root, xml_declaration=True, encoding='UTF-8', standalone=True)

            # Read and modify presentation.xml.rels
            prs_rels_xml = zin.read('ppt/_rels/presentation.xml.rels')
            prs_rels_root = etree.fromstring(prs_rels_xml)
            for rel in list(prs_rels_root):
                if rel.get('Type') == SLIDE_REL_TYPE:
                    prs_rels_root.remove(rel)
            mod_prs_rels_xml = etree.tostring(prs_rels_root, xml_declaration=True, encoding='UTF-8', standalone=True)

            # Read and modify [Content_Types].xml
            ct_xml = zin.read('[Content_Types].xml')
            ct_root = etree.fromstring(ct_xml)
            for override in list(ct_root):
                part_name = override.get('PartName', '')
                if re.match(r'/ppt/slides/slide\d+\.xml$', part_name):
                    ct_root.remove(override)
            mod_ct_xml = etree.tostring(ct_root, xml_declaration=True, encoding='UTF-8', standalone=True)

            with zipfile.ZipFile(tmp_out, 'w', zipfile.ZIP_DEFLATED) as zout:
                skip = {'ppt/presentation.xml', 'ppt/_rels/presentation.xml.rels', '[Content_Types].xml'}
                for item in zin.infolist():
                    if item.filename in to_remove or item.filename in skip:
                        continue
                    zout.writestr(item, zin.read(item.filename))
                zout.writestr('ppt/presentation.xml', mod_prs_xml)
                zout.writestr('ppt/_rels/presentation.xml.rels', mod_prs_rels_xml)
                zout.writestr('[Content_Types].xml', mod_ct_xml)

        os.replace(tmp_out, self.tmp_path)

    def add_slides_from_pptx(self, source_pptx_path):
        """
        Copy ALL slides from source_pptx_path into the output.
        Copies slide XML + all media/chart relationships.
        Returns number of slides added.
        """
        added = 0
        tmp_out = self.tmp_path + '.adding'

        with zipfile.ZipFile(source_pptx_path, 'r') as src_zip:
            src_names = src_zip.namelist()

            # Find all slides in source, sorted by number
            slide_files = sorted(
                [n for n in src_names if re.match(r'ppt/slides/slide\d+\.xml$', n)],
                key=lambda x: int(re.search(r'\d+', x.split('/')[-1]).group())
            )

            if not slide_files:
                return 0

            # Read source content types for media type lookup
            src_ct_xml = src_zip.read('[Content_Types].xml')
            src_ct_root = etree.fromstring(src_ct_xml)
            src_ct_map = {}
            for override in src_ct_root:
                pn = override.get('PartName', '').lstrip('/')
                ct = override.get('ContentType', '')
                if pn and ct:
                    src_ct_map[pn] = ct

            # Extension -> content type fallback
            ext_ct_map = {
                '.png': 'image/png',
                '.jpg': 'image/jpeg',
                '.jpeg': 'image/jpeg',
                '.gif': 'image/gif',
                '.bmp': 'image/bmp',
                '.tiff': 'image/tiff',
                '.emf': 'image/x-emf',
                '.wmf': 'image/x-wmf',
                '.svg': 'image/svg+xml',
            }

            with zipfile.ZipFile(self.tmp_path, 'r') as tgt_zip:
                tgt_names = tgt_zip.namelist()

                # Read target state
                prs_xml = tgt_zip.read('ppt/presentation.xml')
                prs_rels_xml = tgt_zip.read('ppt/_rels/presentation.xml.rels')
                ct_xml = tgt_zip.read('[Content_Types].xml')

                prs_root = etree.fromstring(prs_xml)
                prs_rels_root = etree.fromstring(prs_rels_xml)
                ct_root = etree.fromstring(ct_xml)

                # Count existing slides in target
                existing_slides = [n for n in tgt_names if re.match(r'ppt/slides/slide\d+\.xml$', n)]
                next_slide_num = len(existing_slides) + 1

                # Find max rId in presentation.xml.rels
                existing_rids = set()
                for rel in prs_rels_root:
                    rid = rel.get('Id', '')
                    if rid.startswith('rId'):
                        try:
                            existing_rids.add(int(rid[3:]))
                        except ValueError:
                            pass
                next_prs_rid = max(existing_rids, default=0) + 1

                # Get sldIdLst
                sldIdLst = prs_root.find('.//{%s}sldIdLst' % NS_P)
                if sldIdLst is None:
                    sldIdLst = etree.SubElement(prs_root, '{%s}sldIdLst' % NS_P)

                # Find max sldId
                existing_ids = [int(s.get('id', 255)) for s in sldIdLst]
                next_sld_id = max(existing_ids, default=255) + 1

                # Find a slideLayout in target to use as fallback
                tgt_layouts = sorted(
                    [n for n in tgt_names if re.match(r'ppt/slideLayouts/slideLayout\d+\.xml$', n)],
                    key=lambda x: int(re.search(r'\d+', x.split('/')[-1]).group())
                )
                fallback_layout = tgt_layouts[0] if tgt_layouts else None

                new_files = {}  # path -> bytes

                for slide_path in slide_files:
                    slide_num = next_slide_num
                    next_slide_num += 1

                    new_slide_path = 'ppt/slides/slide{}.xml'.format(slide_num)
                    new_slide_rels_path = 'ppt/slides/_rels/slide{}.xml.rels'.format(slide_num)

                    # Read source slide XML
                    slide_xml_bytes = src_zip.read(slide_path)

                    # Read source slide rels
                    slide_rels_path = slide_path.replace('ppt/slides/', 'ppt/slides/_rels/') + '.rels'
                    src_rels = parse_rels(src_zip, slide_rels_path)

                    # Process each relationship
                    new_rels = []

                    for rel in src_rels:
                        rel_type = rel['type']
                        rel_target = rel['target']
                        rel_mode = rel.get('mode', 'Internal')

                        # External relationships (hyperlinks) — copy as-is
                        if rel_mode == 'External':
                            new_rels.append(rel)
                            continue

                        # slideLayout — point to a layout in the target
                        if SLIDE_LAYOUT_REL_TYPE in rel_type:
                            if fallback_layout:
                                layout_name = fallback_layout.split('/')[-1]
                                new_rels.append({
                                    'id': rel['id'],
                                    'type': rel_type,
                                    'target': '../slideLayouts/' + layout_name,
                                    'mode': 'Internal'
                                })
                            continue

                        # Resolve absolute path in source ZIP
                        if rel_target.startswith('../'):
                            src_abs = 'ppt/' + rel_target[3:]
                        elif rel_target.startswith('/'):
                            src_abs = rel_target.lstrip('/')
                        else:
                            src_abs = 'ppt/slides/' + rel_target

                        # Copy the blob
                        try:
                            blob = src_zip.read(src_abs)
                        except KeyError:
                            print('  [WARN] Rel target not found in source: {}'.format(src_abs), flush=True)
                            continue

                        # Determine content type
                        content_type = src_ct_map.get(src_abs, '')
                        if not content_type:
                            ext = os.path.splitext(src_abs)[1].lower()
                            content_type = ext_ct_map.get(ext, 'application/octet-stream')

                        # Generate unique name for this blob in target
                        ext = os.path.splitext(src_abs)[1]
                        new_blob_name = 'ppt/media/{}{}'.format(uuid.uuid4().hex, ext)
                        new_files[new_blob_name] = blob

                        # Add content type override for new blob
                        etree.SubElement(
                            ct_root,
                            '{%s}Override' % NS_CT,
                            PartName='/' + new_blob_name,
                            ContentType=content_type
                        )

                        # New relative target from slide
                        new_target = '../media/' + os.path.basename(new_blob_name)
                        new_rels.append({
                            'id': rel['id'],
                            'type': rel_type,
                            'target': new_target,
                            'mode': 'Internal'
                        })

                    # Store slide XML and rels
                    new_files[new_slide_path] = slide_xml_bytes
                    new_files[new_slide_rels_path] = build_rels_xml(new_rels)

                    # Add content type for new slide
                    slide_ct = 'application/vnd.openxmlformats-officedocument.presentationml.slide+xml'
                    etree.SubElement(
                        ct_root,
                        '{%s}Override' % NS_CT,
                        PartName='/' + new_slide_path,
                        ContentType=slide_ct
                    )

                    # Add slide to presentation.xml sldIdLst
                    prs_rid = 'rId{}'.format(next_prs_rid)
                    next_prs_rid += 1
                    etree.SubElement(
                        sldIdLst,
                        '{%s}sldId' % NS_P,
                        id=str(next_sld_id),
                        **{'{%s}id' % NS_R: prs_rid}
                    )
                    next_sld_id += 1

                    # Add slide relationship to presentation.xml.rels
                    etree.SubElement(
                        prs_rels_root,
                        '{%s}Relationship' % NS_PR,
                        Id=prs_rid,
                        Type=SLIDE_REL_TYPE,
                        Target='slides/slide{}.xml'.format(slide_num)
                    )

                    added += 1

                # Serialize updated XMLs
                mod_prs_xml = etree.tostring(prs_root, xml_declaration=True, encoding='UTF-8', standalone=True)
                mod_prs_rels_xml = etree.tostring(prs_rels_root, xml_declaration=True, encoding='UTF-8', standalone=True)
                mod_ct_xml = etree.tostring(ct_root, xml_declaration=True, encoding='UTF-8', standalone=True)

                # Write new ZIP
                skip = {'ppt/presentation.xml', 'ppt/_rels/presentation.xml.rels', '[Content_Types].xml'}
                with zipfile.ZipFile(tmp_out, 'w', zipfile.ZIP_DEFLATED) as zout:
                    for item in tgt_zip.infolist():
                        if item.filename not in skip:
                            zout.writestr(item, tgt_zip.read(item.filename))
                    zout.writestr('ppt/presentation.xml', mod_prs_xml)
                    zout.writestr('ppt/_rels/presentation.xml.rels', mod_prs_rels_xml)
                    zout.writestr('[Content_Types].xml', mod_ct_xml)
                    for path, data in new_files.items():
                        zout.writestr(path, data)

        os.replace(tmp_out, self.tmp_path)
        self.slide_count += added
        return added

    def save(self, output_path):
        shutil.copy2(self.tmp_path, output_path)
        try:
            os.remove(self.tmp_path)
        except Exception:
            pass

    def __del__(self):
        if hasattr(self, 'tmp_path') and self.tmp_path and os.path.exists(self.tmp_path):
            try:
                os.remove(self.tmp_path)
            except Exception:
                pass


# ─────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────
def main():
    try:
        if len(sys.argv) < 2:
            raise Exception('Usage: python assemble_engine.py <payload.json>')

        payload_path = sys.argv[1]
        with open(payload_path, 'r', encoding='utf-8') as f:
            data = json.load(f)

        LIBRARY_ROOT = data['library_root']
        OUTPUT_DIR = data['output_dir']
        presentation_type_name = data.get('presentation_type', 'Feasibility Study')
        sections = data.get('sections', [])
        form_data = data.get('formData', {})
        plots = data.get('plots', [])

        os.makedirs(OUTPUT_DIR, exist_ok=True)

        # 1. Initialize assembler
        root_template = os.path.join(LIBRARY_ROOT, 'RootTemplate.pptx')
        assembler = PptxAssembler(template_path=root_template if os.path.exists(root_template) else None)
        if os.path.exists(root_template):
            print('[INFO] Initialized from RootTemplate.pptx', flush=True)
        else:
            print('[WARN] RootTemplate.pptx not found, using blank presentation', flush=True)

        # 2. Sort sections by order
        sections.sort(key=lambda s: s.get('order', 999))

        library_base = os.path.join(LIBRARY_ROOT, presentation_type_name)
        total_slides_added = 0
        skipped_sections = []
        added_sections = []

        # 3. Process each section
        for section in sections:
            s_name = section.get('name', 'Unknown')
            s_folder = section.get('folderPath', s_name)
            s_varying = section.get('isVarying', False)
            s_criteria_keys = section.get('varyingCriteria', ['City', 'Asset Type', 'Category', 'Specifications'])
            s_filename = section.get('filename')

            section_dir = os.path.join(library_base, s_folder)
            print('[SECTION] {} | Varying: {} | Dir: {}'.format(s_name, s_varying, section_dir), flush=True)

            if not os.path.isdir(section_dir):
                skipped_sections.append('{} - folder not found: {}'.format(s_name, s_folder))
                print('  [SKIP] Folder not found', flush=True)
                continue

            # NON-VARYING SECTION
            if not s_varying:
                source_path = find_pptx_in_folder(section_dir, s_filename)
                if not source_path:
                    skipped_sections.append('{} - no PPTX file in folder'.format(s_name))
                    print('  [SKIP] No PPTX in folder', flush=True)
                    continue

                print('  [FIXED] Loading: {}'.format(os.path.basename(source_path)), flush=True)
                try:
                    added = assembler.add_slides_from_pptx(source_path)
                    total_slides_added += added
                    added_sections.append('{} ({} slides)'.format(s_name, added))
                    print('  [OK] Added {} slides'.format(added), flush=True)
                except Exception as e:
                    skipped_sections.append('{} - error: {}'.format(s_name, str(e)))
                    print('  [ERROR] {}'.format(str(e)), flush=True)

            # VARYING SECTION
            else:
                if not plots:
                    skipped_sections.append('{} - no plots provided'.format(s_name))
                    print('  [SKIP] No plots for varying section', flush=True)
                    continue

                unique_criteria_list = deduplicate_plots(plots, s_criteria_keys)
                print('  [VARYING] {} plots -> {} unique combinations'.format(len(plots), len(unique_criteria_list)), flush=True)

                section_slides_added = 0
                for criteria in unique_criteria_list:
                    values = []
                    for key in s_criteria_keys:
                        val = next((v for ck, v in criteria.items() if ck.lower() == key.lower()), None)
                        if val:
                            values.append(val)

                    if not values:
                        print('  [SKIP] No criteria values found: {}'.format(criteria), flush=True)
                        continue

                    # Try all naming conventions (underscore, plus-space, partial)
                    source_path, matched_name = find_pptx_for_criteria(section_dir, values)

                    # Log what we tried
                    tried = build_varying_filenames(values)
                    print('  [LOOKUP] Tried: {} -> {}'.format(
                        tried, 'FOUND: ' + matched_name if source_path else 'NOT FOUND'
                    ), flush=True)

                    if not source_path:
                        skipped_sections.append('{} - no match for: {}'.format(s_name, tried[0] if tried else str(values)))
                        continue

                    try:
                        added = assembler.add_slides_from_pptx(source_path)
                        total_slides_added += added
                        section_slides_added += added
                        print('  [OK] Added {} slides from {}'.format(added, os.path.basename(source_path)), flush=True)
                    except Exception as e:
                        skipped_sections.append('{} [{}] - error: {}'.format(s_name, matched_name, str(e)))
                        print('  [ERROR] {}'.format(str(e)), flush=True)

                if section_slides_added > 0:
                    added_sections.append('{} ({} slides)'.format(s_name, section_slides_added))

        # 4. Save output
        safe_title = re.sub(r'[^a-zA-Z0-9_\-]', '_', str(form_data.get('title', form_data.get('projectTitle', 'Presentation'))))
        safe_title = safe_title[:50]
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        output_filename = '{}_{}_{}.pptx'.format(safe_title, timestamp, uuid.uuid4().hex[:6])
        output_path = os.path.join(OUTPUT_DIR, output_filename)

        assembler.save(output_path)

        print(json.dumps({
            'success': True,
            'fileName': output_filename,
            'filePath': output_path,
            'slideCount': total_slides_added,
            'addedSections': added_sections,
            'skippedSections': skipped_sections
        }), flush=True)

    except Exception as e:
        import traceback
        tb = traceback.format_exc()
        print('[FATAL]\n' + tb, file=sys.stderr, flush=True)
        print(json.dumps({'success': False, 'error': str(e), 'traceback': tb}), flush=True)
        sys.exit(1)


if __name__ == '__main__':
    main()
