
import sys
import os
import json
import re
from datetime import datetime

# Try to import python-pptx, handle missing dependency
try:
    from pptx import Presentation
    from pptx.util import Inches, Emu
    from copy import deepcopy
except ImportError:
    print(json.dumps({"success": False, "error": "Missing python-pptx library. Auto-installation should have happened."}))
    sys.exit(1)

# FORCE CORRECT DIMENSIONS (20 x 11.2 inches)
WIDTH_INCHES = 20
HEIGHT_INCHES = 11.2

def set_dimensions(prs):
    prs.slide_width = Inches(WIDTH_INCHES)
    prs.slide_height = Inches(HEIGHT_INCHES)

def copy_slide_content(source_slide, target_prs):
    """Deep copy content from source slide to new slide in target"""
    try:
        # Create a blank slide (using last layout to be safe, or blank)
        # We try to find a blank layout
        layout = target_prs.slide_layouts[6] if len(target_prs.slide_layouts) > 6 else target_prs.slide_layouts[0]
        
        target_slide = target_prs.slides.add_slide(layout)
        
        # Copy all shapes with full fidelity (charts, tables, images)
        for shape in source_slide.shapes:
            el = shape.element
            newel = deepcopy(el)
            target_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
            
        return True
    except Exception as e:
        # Log warning but continue
        return False

def resolve_placeholders(prs, data):
    replacements = {
        '{{PROJECT_NAME}}': data.get('title', ''),
        '{{Title}}': data.get('title', ''),
        '{{Subtitle}}': data.get('subtitle', ''),
        '{{CLIENT_NAME}}': data.get('clientName', ''),
        '{{DATE}}': datetime.now().strftime('%B %d, %Y'),
        '{{CITY}}': data.get('city', ''),
        '{{ASSET_TYPE}}': data.get('assetType', ''),
        '{{CATEGORY}}': data.get('category', ''),
        '{{SPECIFICATIONS}}': data.get('specifications', '')
    }
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                try:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text = run.text
                            for key, val in replacements.items():
                                if key in text:
                                    text = text.replace(key, str(val))
                            run.text = text
                except:
                    pass

def main():
    try:
        if len(sys.argv) < 2:
            raise Exception("Missing payload file argument")

        payload_path = sys.argv[1]
        
        with open(payload_path, 'r', encoding='utf-8') as f:
            data = json.load(f)

        # Config
        LIBRARY_ROOT = data.get('library_root')
        OUTPUT_DIR = data.get('output_dir')
        
        p_type_name = data.get('presentation_type', 'Feasibility Study')
        p_config = data.get('presentation_config', {})
        form_data = data.get('formData', {})
        plots = data.get('plots', [])

        # 1. Initialize Master Presentation
        # Check for RootTemplate
        root_template_path = os.path.join(LIBRARY_ROOT, 'RootTemplate.pptx')
        if os.path.exists(root_template_path):
            prs = Presentation(root_template_path)
            # Remove existing slides from template to start clean, but keep master layouts
            xml_slides = prs.slides._sldIdLst
            slides = list(xml_slides)
            for s in slides:
                xml_slides.remove(s)
        else:
            prs = Presentation()

        set_dimensions(prs)

        # 2. Prepare Contexts (Plots + Global)
        # "Plots enabled? Then each plot has own answers."
        # "Also need to handle global context."
        context_list = [p.get('criteria', {}) for p in plots] if plots else [form_data]
        # Filter empty contexts
        context_list = [ctx for ctx in context_list if ctx]
        if not context_list: 
            context_list = [form_data] # Fallback to global form data

        # 3. Process Sections
        sections = p_config.get('sections', [])
        # Ensure correct order
        sections.sort(key=lambda x: x.get('order', 0))

        skipped_sections = []
        total_slides = 0

        library_base = os.path.join(LIBRARY_ROOT, p_type_name)

        for section in sections:
            s_name = section.get('name')
            s_folder = section.get('folderPath', s_name)
            s_varying = section.get('isVarying', False)
            s_criteria = section.get('varyingCriteria', ['City', 'Asset Type'])
            
            section_dir = os.path.join(library_base, s_folder)
            
            if not os.path.exists(section_dir):
                skipped_sections.append(f"{s_name} (Folder missing)")
                continue

            # --- UNVARYING / FIXED SECTION ---
            if not s_varying:
                # Logic: Take the first PPTX found in folder
                files = [f for f in os.listdir(section_dir) if f.lower().endswith('.pptx') and not f.startswith('~$')]
                if not files:
                    skipped_sections.append(f"{s_name} (No files)")
                    continue
                
                # We need to add this slide ONCE (global)
                source_path = os.path.join(section_dir, files[0])
                if os.path.exists(source_path):
                   source_prs = Presentation(source_path)
                   for slide in source_prs.slides:
                       if copy_slide_content(slide, prs):
                           total_slides += 1
            
            # --- VARYING SECTION (The core requirement) ---
            else:
                # Logic: For each "Plot" (context), find the matching PPTX.
                # Overlap Deduction: If multiple plots point to SAME file, only add ONCE per block?
                # The requirement says: "If two plots generate the same PPTX... then there will be ONE set of slides"
                # So we deduce overlap PER SECTION.
                
                matched_files = set()
                files_to_add = [] # List of unique paths to add for this section

                for ctx in context_list:
                    # Build filename: "Value1 + Value2.pptx"
                    parts = []
                    for crit in s_criteria:
                        # Find key case-insensitive
                        val = next((v for k,v in ctx.items() if k.lower() == crit.lower()), None)
                        if val:
                            parts.append(val)
                    
                    if not parts:
                         continue

                    # Try permutations or precise matching? 
                    # Requirement implies: "Riyadh + Residential + Apartments + Luxury"
                    # User likely has filenames exactly like this.
                    filename_base = " + ".join(parts)
                    filename = f"{filename_base}.pptx"
                    file_path = os.path.join(section_dir, filename)

                    # Check existence (Case insensitive check strictly speaking required on Linux, Windows is forgiving)
                    # We'll trust os.path.exists on Windows
                    if os.path.exists(file_path):
                        if file_path not in matched_files:
                            matched_files.add(file_path)
                            files_to_add.append(file_path)
                    else:
                        # Try finding file ignoring case
                        all_files = os.listdir(section_dir)
                        found = False
                        for f in all_files:
                            if f.lower() == filename.lower():
                                real_path = os.path.join(section_dir, f)
                                if real_path not in matched_files:
                                    matched_files.add(real_path)
                                    files_to_add.append(real_path)
                                found = True
                                break
                
                # Add the Unique found files
                for fp in files_to_add:
                    src = Presentation(fp)
                    for slide in src.slides:
                         if copy_slide_content(slide, prs):
                             total_slides += 1
                
                if not files_to_add:
                    skipped_sections.append(f"{s_name} (No matches found)")

        # 4. Final Polish
        resolve_placeholders(prs, form_data)
        
        # Save
        safe_title = re.sub(r'[^a-zA-Z0-9]', '_', form_data.get('title', 'Presentation'))
        filename = f"{safe_title}_{uuidv4()[:8]}.pptx"
        output_path = os.path.join(OUTPUT_DIR, filename)
        
        prs.save(output_path)
        
        print(json.dumps({
            "success": True,
            "fileName": filename,
            "filePath": output_path,
            "slideCount": total_slides,
            "skippedSections": skipped_sections
        }))

    except Exception as e:
        import traceback
        traceback.print_exc()
        print(json.dumps({"success": False, "error": str(e)}))

if __name__ == '__main__':
    main()
